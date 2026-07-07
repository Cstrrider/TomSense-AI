"""File-upload handling: classify, extract, persist.

Storage layout:
    /data/uploads/<uuid>.bin     — raw bytes for images
    /data/uploads/<uuid>.txt     — original bytes for text (so frontend can re-fetch)
    /data/uploads/<uuid>.pdf     — original PDF (text already extracted into DB)

For text and PDF the extracted body is also stored as a truncated excerpt in
DB, so the chat path can prepend it without re-reading the file from disk.
"""

import base64
import io
import os
import uuid
from typing import Optional

UPLOADS_DIR = os.getenv("UPLOADS_DIR", "/data/uploads")

# Caps — phone photos routinely come in at 10-15 MB, so we accept generously
# at ingest and let `normalize_image` downscale before persisting.
MAX_IMAGE_BYTES = 30 * 1024 * 1024    # 30 MB raw
MAX_TEXT_BYTES  = 2 * 1024 * 1024     # 2 MB
MAX_PDF_BYTES   = 32 * 1024 * 1024    # 32 MB
MAX_DOC_BYTES   = 20 * 1024 * 1024    # 20 MB (docx/xlsx/pptx)
MAX_AUDIO_BYTES = 25 * 1024 * 1024    # 25 MB (whisper handles ~30 min mp3)
MAX_EXCERPT_CHARS = 50_000

# Downscale target: longest edge in pixels. Gemma 4 vision doesn't benefit
# from anything larger and base64 inline-image tokens balloon fast.
IMAGE_MAX_EDGE = 1600
IMAGE_JPEG_QUALITY = 85

IMAGE_MIMES = {"image/jpeg", "image/png", "image/webp", "image/gif"}
TEXT_MIMES_PREFIX = ("text/",)
TEXT_EXTRA_MIMES = {
    "application/json", "application/xml", "application/x-yaml",
    "application/javascript", "application/typescript",
}
PDF_MIMES = {"application/pdf"}
AUDIO_EXTS = {"mp3", "wav", "m4a", "ogg", "oga", "webm", "flac", "aac", "opus"}
DOC_EXTS = {"docx", "xlsx", "xlsm", "pptx"}

# Kinds whose content is text-in-DB (text_excerpt) — these flow through the
# chat prepend path and are RAG-indexable.
TEXTUAL_KINDS = ("text", "pdf", "doc", "audio")


def ensure_dir() -> None:
    os.makedirs(UPLOADS_DIR, exist_ok=True)


def classify(mime: str, filename: str) -> Optional[str]:
    m = (mime or "").lower()
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if m in IMAGE_MIMES or ext in {"jpg", "jpeg", "png", "webp", "gif"}:
        return "image"
    if m in PDF_MIMES or ext == "pdf":
        return "pdf"
    if m.startswith("audio/") or ext in AUDIO_EXTS:
        return "audio"
    if ext in DOC_EXTS:
        return "doc"
    if m.startswith(TEXT_MIMES_PREFIX) or m in TEXT_EXTRA_MIMES or ext in {
        "txt", "md", "markdown", "rst", "csv", "tsv", "log",
        "py", "js", "ts", "tsx", "jsx", "json", "yaml", "yml",
        "html", "css", "sh", "bash", "rb", "go", "rs", "java",
        "kt", "swift", "c", "h", "cpp", "hpp", "sql", "toml", "ini",
    }:
        return "text"
    return None


def _path_for(upload_id: str, ext: str) -> str:
    return os.path.join(UPLOADS_DIR, f"{upload_id}.{ext}")


def normalize_image(raw: bytes, mime: str) -> tuple[bytes, str, str]:
    """Auto-rotate via EXIF, downscale to IMAGE_MAX_EDGE, re-encode efficiently.

    Returns (new_bytes, new_mime, new_ext). Falls back to the original bytes
    if Pillow can't decode the input.
    """
    try:
        from PIL import Image, ImageOps
    except Exception:
        return raw, mime, _ext_from_mime(mime)

    try:
        img = Image.open(io.BytesIO(raw))
        img = ImageOps.exif_transpose(img)
        # Downscale if either dimension exceeds the cap
        w, h = img.size
        longest = max(w, h)
        if longest > IMAGE_MAX_EDGE:
            scale = IMAGE_MAX_EDGE / longest
            img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

        # Re-encode. PNG with alpha → keep PNG, else go JPEG for size.
        has_alpha = img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info)
        out = io.BytesIO()
        if has_alpha:
            img.save(out, format="PNG", optimize=True)
            new_bytes = out.getvalue()
            # If JPEG-with-flat-bg would be smaller, prefer it
            return new_bytes, "image/png", "png"
        else:
            img = img.convert("RGB")
            img.save(out, format="JPEG", quality=IMAGE_JPEG_QUALITY, optimize=True, progressive=True)
            new_bytes = out.getvalue()
            return new_bytes, "image/jpeg", "jpg"
    except Exception:
        return raw, mime, _ext_from_mime(mime)


def _ext_from_mime(mime: str) -> str:
    m = (mime or "").lower()
    if m == "image/jpeg" or m == "image/jpg":
        return "jpg"
    if m == "image/png":
        return "png"
    if m == "image/webp":
        return "webp"
    if m == "image/gif":
        return "gif"
    return "bin"


def extract_pdf_text(raw: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return ""
    try:
        reader = PdfReader(io.BytesIO(raw))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n\n".join(p for p in parts if p)
    except Exception:
        return ""


def extract_docx_text(raw: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""


def extract_xlsx_text(raw: bytes) -> str:
    """Sheets rendered as markdown-ish pipe tables, capped per sheet so a
    100k-row export can't produce a gigabyte of text."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 500:
                    parts.append(f"… [{ws.max_row - 500} more rows truncated]")
                    break
                parts.append(" | ".join("" if v is None else str(v) for v in row))
        wb.close()
        return "\n".join(parts)
    except Exception:
        return ""


def extract_pptx_text(raw: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        return ""


def extract_doc_text(raw: bytes, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    if ext == "docx":
        return extract_docx_text(raw)
    if ext in ("xlsx", "xlsm"):
        return extract_xlsx_text(raw)
    if ext == "pptx":
        return extract_pptx_text(raw)
    return ""


def render_pdf_pages(raw: bytes, max_pages: int = 4, dpi: int = 120) -> list[bytes]:
    """Render the first pages of a PDF to JPEG — OCR fallback for scanned
    documents where pypdf finds no text layer."""
    try:
        import fitz  # PyMuPDF
    except Exception:
        return []
    out: list[bytes] = []
    try:
        doc = fitz.open(stream=raw, filetype="pdf")
        for page in doc[:max_pages]:
            pix = page.get_pixmap(dpi=dpi)
            out.append(pix.tobytes("jpeg"))
        doc.close()
    except Exception:
        return []
    return out


def process_upload(filename: str, mime: str, raw: bytes) -> dict:
    """Validate, extract, and persist bytes to disk.

    Returns a metadata dict suitable for `db.insert_upload(...)`. The DB write
    itself happens in the endpoint.
    """
    ensure_dir()
    kind = classify(mime, filename)
    if kind is None:
        raise ValueError(f"unsupported file type: {mime} ({filename})")

    size = len(raw)
    if kind == "image" and size > MAX_IMAGE_BYTES:
        raise ValueError(f"image too large ({size} > {MAX_IMAGE_BYTES})")
    if kind == "text" and size > MAX_TEXT_BYTES:
        raise ValueError(f"text too large ({size} > {MAX_TEXT_BYTES})")
    if kind == "pdf"  and size > MAX_PDF_BYTES:
        raise ValueError(f"pdf too large ({size} > {MAX_PDF_BYTES})")
    if kind == "doc" and size > MAX_DOC_BYTES:
        raise ValueError(f"document too large ({size} > {MAX_DOC_BYTES})")
    if kind == "audio" and size > MAX_AUDIO_BYTES:
        raise ValueError(f"audio too large ({size} > {MAX_AUDIO_BYTES})")

    upload_id = str(uuid.uuid4())
    excerpt: Optional[str] = None
    storage_path: Optional[str] = None

    if kind == "image":
        new_bytes, new_mime, ext = normalize_image(raw, mime)
        mime = new_mime
        size = len(new_bytes)
        storage_path = _path_for(upload_id, ext)
        with open(storage_path, "wb") as f:
            f.write(new_bytes)
    elif kind == "text":
        try:
            text = raw.decode("utf-8", errors="replace")
        except Exception:
            text = ""
        excerpt = text[:MAX_EXCERPT_CHARS]
        storage_path = _path_for(upload_id, "txt")
        with open(storage_path, "wb") as f:
            f.write(raw)
    elif kind == "pdf":
        excerpt = extract_pdf_text(raw)[:MAX_EXCERPT_CHARS]
        storage_path = _path_for(upload_id, "pdf")
        with open(storage_path, "wb") as f:
            f.write(raw)
    elif kind == "doc":
        excerpt = extract_doc_text(raw, filename)[:MAX_EXCERPT_CHARS]
        ext = filename.lower().rsplit(".", 1)[-1]
        storage_path = _path_for(upload_id, ext)
        with open(storage_path, "wb") as f:
            f.write(raw)
    elif kind == "audio":
        # excerpt (the transcript) is filled in by the /uploads endpoint —
        # transcription is an async Whisper call, this function is sync.
        ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else "bin"
        storage_path = _path_for(upload_id, ext)
        with open(storage_path, "wb") as f:
            f.write(raw)

    return {
        "upload_id": upload_id,
        "kind": kind,
        "filename": filename,
        "mime": mime,
        "size_bytes": size,
        "text_excerpt": excerpt,
        "storage_path": storage_path,
    }


async def read_full_text(meta: dict) -> str:
    """Return the FULL text body of an upload (for RAG indexing).

    For text uploads, decodes the file on disk. For PDFs, re-runs pypdf so we
    can index the whole document, not just the first MAX_EXCERPT_CHARS that
    landed in `text_excerpt`.
    """
    path = meta.get("storage_path")
    kind = meta.get("kind")
    if not path or not kind:
        return ""
    try:
        with open(path, "rb") as f:
            raw = f.read()
    except OSError:
        return ""
    if kind == "text":
        try:
            return raw.decode("utf-8", errors="replace")
        except Exception:
            return ""
    if kind == "pdf":
        return extract_pdf_text(raw)
    if kind == "doc":
        return extract_doc_text(raw, meta.get("filename") or "")
    if kind == "audio":
        # Transcript was produced once at ingest (whisper); it IS the text.
        return meta.get("text_excerpt") or ""
    return ""


def image_data_url(meta: dict) -> Optional[str]:
    """Read the image bytes back and return a data: URL ready for vision input."""
    if meta.get("kind") != "image" or not meta.get("storage_path"):
        return None
    try:
        with open(meta["storage_path"], "rb") as f:
            raw = f.read()
    except OSError:
        return None
    b64 = base64.b64encode(raw).decode("ascii")
    mime = meta.get("mime") or "image/jpeg"
    return f"data:{mime};base64,{b64}"
